import os, json, io, sqlite3, uuid, re
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ────────────────────────────────────────────────────
BASE_DIR      = os.path.dirname(os.path.abspath(__file__))   # .../server
ROOT_DIR      = os.path.dirname(BASE_DIR)                     # .../lectureai
TEMPLATE_PATH = os.path.join(ROOT_DIR, "client", "templates", "index.html")
DB_PATH       = "/tmp/lectureai.db"

app = Flask(__name__, static_folder="static", template_folder="templates")
CORS(app)
client = Groq(api_key=os.environ.get("GROQ_API_KEY"))

# ══════════════════════════════════════════════════════════════
#  DATABASE
# ══════════════════════════════════════════════════════════════
def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    try:
        conn = get_db()
        conn.execute("""CREATE TABLE IF NOT EXISTS classes (
            code TEXT PRIMARY KEY, teacher_email TEXT, teacher_name TEXT,
            topic TEXT, level TEXT, data TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS assignments (
            id TEXT PRIMARY KEY, class_code TEXT, teacher_email TEXT,
            title TEXT, description TEXT, due_date TEXT,
            max_score INTEGER DEFAULT 100, created_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS submissions (
            id TEXT PRIMARY KEY, assignment_id TEXT, class_code TEXT,
            student_email TEXT, student_name TEXT, content TEXT,
            submitted_at TEXT, score INTEGER DEFAULT -1, feedback TEXT DEFAULT '')""")
        conn.execute("""CREATE TABLE IF NOT EXISTS tests (
            id TEXT PRIMARY KEY, class_code TEXT, teacher_email TEXT,
            title TEXT, questions TEXT, time_limit INTEGER DEFAULT 0, created_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS test_submissions (
            id TEXT PRIMARY KEY, test_id TEXT, class_code TEXT,
            student_email TEXT, student_name TEXT, answers TEXT,
            score INTEGER DEFAULT 0, total INTEGER DEFAULT 0, submitted_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS attendance (
            id TEXT PRIMARY KEY, class_code TEXT, teacher_email TEXT,
            student_name TEXT, session_date TEXT,
            present INTEGER DEFAULT 1, created_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS discussions (
            id TEXT PRIMARY KEY, class_code TEXT, student_name TEXT,
            student_email TEXT, question TEXT, reply TEXT DEFAULT '',
            replied_by TEXT DEFAULT '', created_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS reactions (
            id TEXT PRIMARY KEY, class_code TEXT, student_email TEXT,
            student_name TEXT, reaction TEXT, created_at TEXT)""")
        conn.execute("""CREATE TABLE IF NOT EXISTS lecture_library (
            id TEXT PRIMARY KEY,
            teacher_email TEXT NOT NULL,
            teacher_name TEXT,
            title TEXT NOT NULL,
            topic TEXT NOT NULL,
            subject TEXT DEFAULT '',
            level TEXT DEFAULT 'Intermediate',
            institution TEXT DEFAULT '',
            notes TEXT NOT NULL,
            class_code TEXT DEFAULT '',
            is_public INTEGER DEFAULT 1,
            view_count INTEGER DEFAULT 0,
            saved_at TEXT NOT NULL,
            year TEXT DEFAULT '')""")
        conn.commit()
        conn.close()
    except Exception as e:
        print("DB init error:", e)

init_db()

def ask_groq(prompt, max_tokens=1500):
    r = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=max_tokens,
        temperature=0.7
    )
    return r.choices[0].message.content.strip()

# ══════════════════════════════════════════════════════════════
#  CORE ROUTES
# ══════════════════════════════════════════════════════════════
@app.route("/")
def index():
    with open(TEMPLATE_PATH, encoding="utf-8") as f:
        return f.read()

@app.route("/ping")
def ping():
    return jsonify({"status": "ok", "version": "2.1"})

@app.route("/health")
def health():
    return jsonify({"status": "ok", "db": os.path.exists(DB_PATH), "template": os.path.exists(TEMPLATE_PATH)})

# ══════════════════════════════════════════════════════════════
#  CLASS MANAGEMENT
# ══════════════════════════════════════════════════════════════
@app.route("/save_class", methods=["POST"])
def save_class():
    try:
        d = request.json
        code = d.get("code", "").upper().strip()
        if not code:
            return jsonify({"success": False, "error": "No code provided"})
        init_db()
        conn = get_db()
        conn.execute("""INSERT INTO classes (code,teacher_email,teacher_name,topic,level,data)
            VALUES(?,?,?,?,?,?) ON CONFLICT(code) DO UPDATE SET
            teacher_email=excluded.teacher_email, teacher_name=excluded.teacher_name,
            topic=excluded.topic, level=excluded.level, data=excluded.data""",
            (code, d.get("teacherEmail",""), d.get("teacherName",""),
             d.get("topic",""), d.get("level",""), json.dumps(d)))
        conn.commit()
        row = conn.execute("SELECT code FROM classes WHERE code=?", (code,)).fetchone()
        conn.close()
        if row:
            print(f"[LectureAI] Class saved: {code}")
            return jsonify({"success": True, "code": code})
        return jsonify({"success": False, "error": "Save failed verification"})
    except Exception as e:
        print(f"[LectureAI] save_class error: {e}")
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_class", methods=["POST"])
def get_class():
    try:
        d = request.json
        code = d.get("code", "").upper().strip()
        if not code:
            return jsonify({"success": False, "error": "No code provided"})
        init_db()
        conn = get_db()
        row = conn.execute("SELECT data FROM classes WHERE code=?", (code,)).fetchone()
        count = conn.execute("SELECT COUNT(*) as cnt FROM classes").fetchone()
        conn.close()
        if row:
            print(f"[LectureAI] Class found: {code}")
            return jsonify({"success": True, "class": json.loads(row["data"])})
        print(f"[LectureAI] Class NOT found: '{code}' (total: {count['cnt']})")
        return jsonify({"success": False, "error": "Code not found"})
    except Exception as e:
        print(f"[LectureAI] get_class error: {e}")
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  ASSIGNMENTS
# ══════════════════════════════════════════════════════════════
@app.route("/create_assignment", methods=["POST"])
def create_assignment():
    try:
        d = request.json
        aid = str(uuid.uuid4())
        init_db()
        conn = get_db()
        conn.execute("""INSERT INTO assignments
            (id,class_code,teacher_email,title,description,due_date,max_score,created_at)
            VALUES(?,?,?,?,?,?,?,datetime('now'))""",
            (aid, d.get("classCode",""), d.get("teacherEmail",""), d.get("title",""),
             d.get("description",""), d.get("dueDate",""), int(d.get("maxScore",100))))
        conn.commit()
        conn.close()
        return jsonify({"success": True, "id": aid})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_assignments", methods=["POST"])
def get_assignments():
    try:
        d = request.json
        code = d.get("classCode","").upper().strip()
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM assignments WHERE class_code=? ORDER BY created_at DESC", (code,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "assignments": [dict(r) for r in rows]})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/delete_assignment", methods=["POST"])
def delete_assignment():
    try:
        d = request.json
        aid = d.get("id","")
        init_db()
        conn = get_db()
        conn.execute("DELETE FROM assignments WHERE id=?", (aid,))
        conn.execute("DELETE FROM submissions WHERE assignment_id=?", (aid,))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/submit_assignment", methods=["POST"])
def submit_assignment():
    try:
        d = request.json
        init_db()
        conn = get_db()
        existing = conn.execute(
            "SELECT id FROM submissions WHERE assignment_id=? AND student_email=?",
            (d.get("assignmentId",""), d.get("studentEmail",""))
        ).fetchone()
        if existing:
            conn.execute(
                "UPDATE submissions SET content=?,submitted_at=datetime('now') WHERE id=?",
                (d.get("content",""), existing["id"])
            )
        else:
            conn.execute("""INSERT INTO submissions
                (id,assignment_id,class_code,student_email,student_name,content,submitted_at)
                VALUES(?,?,?,?,?,?,datetime('now'))""",
                (str(uuid.uuid4()), d.get("assignmentId",""), d.get("classCode",""),
                 d.get("studentEmail",""), d.get("studentName",""), d.get("content","")))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_submissions", methods=["POST"])
def get_submissions():
    try:
        d = request.json
        aid = d.get("assignmentId","")
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM submissions WHERE assignment_id=? ORDER BY submitted_at DESC", (aid,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "submissions": [dict(r) for r in rows]})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/grade_submission", methods=["POST"])
def grade_submission():
    try:
        d = request.json
        init_db()
        conn = get_db()
        conn.execute(
            "UPDATE submissions SET score=?,feedback=? WHERE id=?",
            (int(d.get("score",0)), d.get("feedback",""), d.get("submissionId",""))
        )
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_my_submission", methods=["POST"])
def get_my_submission():
    try:
        d = request.json
        init_db()
        conn = get_db()
        row = conn.execute(
            "SELECT * FROM submissions WHERE assignment_id=? AND student_email=?",
            (d.get("assignmentId",""), d.get("studentEmail",""))
        ).fetchone()
        conn.close()
        return jsonify({"success": True, "submission": dict(row) if row else None})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  TESTS
# ══════════════════════════════════════════════════════════════
@app.route("/create_test", methods=["POST"])
def create_test():
    try:
        d = request.json
        tid = str(uuid.uuid4())
        init_db()
        conn = get_db()
        conn.execute("""INSERT INTO tests
            (id,class_code,teacher_email,title,questions,time_limit,created_at)
            VALUES(?,?,?,?,?,?,datetime('now'))""",
            (tid, d.get("classCode",""), d.get("teacherEmail",""), d.get("title",""),
             json.dumps(d.get("questions",[])), int(d.get("timeLimit",0))))
        conn.commit()
        conn.close()
        return jsonify({"success": True, "id": tid})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_tests", methods=["POST"])
def get_tests():
    try:
        d = request.json
        code = d.get("classCode","").upper().strip()
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM tests WHERE class_code=? ORDER BY created_at DESC", (code,)
        ).fetchall()
        result = []
        for r in rows:
            rd = dict(r)
            rd["questions"] = json.loads(rd["questions"])
            result.append(rd)
        conn.close()
        return jsonify({"success": True, "tests": result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/delete_test", methods=["POST"])
def delete_test():
    try:
        d = request.json
        tid = d.get("id","")
        init_db()
        conn = get_db()
        conn.execute("DELETE FROM tests WHERE id=?", (tid,))
        conn.execute("DELETE FROM test_submissions WHERE test_id=?", (tid,))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/submit_test", methods=["POST"])
def submit_test():
    try:
        d = request.json
        answers = d.get("answers", {})
        questions = d.get("questions", [])
        score = sum(
            1 for i, q in enumerate(questions)
            if str(i) in answers and int(answers[str(i)]) == int(q.get("ans", -1))
        )
        init_db()
        conn = get_db()
        existing = conn.execute(
            "SELECT id FROM test_submissions WHERE test_id=? AND student_email=?",
            (d.get("testId",""), d.get("studentEmail",""))
        ).fetchone()
        if existing:
            return jsonify({"success": False, "error": "Already submitted"})
        conn.execute("""INSERT INTO test_submissions
            (id,test_id,class_code,student_email,student_name,answers,score,total,submitted_at)
            VALUES(?,?,?,?,?,?,?,?,datetime('now'))""",
            (str(uuid.uuid4()), d.get("testId",""), d.get("classCode",""),
             d.get("studentEmail",""), d.get("studentName",""),
             json.dumps(answers), score, len(questions)))
        conn.commit()
        conn.close()
        return jsonify({"success": True, "score": score, "total": len(questions)})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_test_results", methods=["POST"])
def get_test_results():
    try:
        d = request.json
        tid = d.get("testId","")
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM test_submissions WHERE test_id=? ORDER BY submitted_at DESC", (tid,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "results": [dict(r) for r in rows]})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_my_test_result", methods=["POST"])
def get_my_test_result():
    try:
        d = request.json
        init_db()
        conn = get_db()
        row = conn.execute(
            "SELECT * FROM test_submissions WHERE test_id=? AND student_email=?",
            (d.get("testId",""), d.get("studentEmail",""))
        ).fetchone()
        conn.close()
        return jsonify({"success": True, "result": dict(row) if row else None})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  AI: ICAP LECTURE NOTES
# ══════════════════════════════════════════════════════════════
@app.route("/generate_notes", methods=["POST"])
def generate_notes():
    try:
        d = request.json
        topic = d.get("topic","")
        level = d.get("level","Intermediate")
        duration = d.get("duration", 75)
        objectives = d.get("objectives","")
        style = d.get("style","Lecture-based")
        language = d.get("language","English")

        p = f"""You are an expert university lecturer writing comprehensive, well-structured lecture notes on "{topic}" for {level}-level students. Duration: {duration} minutes. Style: {style}.

IMPORTANT: Write the ENTIRE response in {language}.

Learning objectives:
{objectives if objectives else f"Cover {topic} comprehensively for {level}-level learners."}

Structure the notes in exactly this order. Use the ICAP tags exactly as shown — they are used for colour coding only and must be kept:

[PASSIVE] 1. INTRODUCTION AND CONTEXT
Why this topic matters. Historical background. Real-world relevance. A compelling hook. At least 3 detailed paragraphs.

[PASSIVE] 2. CORE DEFINITIONS AND TERMINOLOGY
Define every key term precisely with examples. At least 6-8 terms. Use **Term**: definition format for each.

[ACTIVE] 3. CORE CONCEPT EXPLANATIONS
Deep dive into each major concept. Step-by-step reasoning. Multiple representations. At least 3-4 major concepts, each with 2+ paragraphs and bullet points.

[ACTIVE] 4. WORKED EXAMPLES
Use "Example:" to start each worked example. At least 3 fully worked examples of increasing difficulty. Show every step and explain WHY.

[CONSTRUCTIVE] 5. CRITICAL THINKING & ANALYSIS
Open-ended analysis questions. "What if" scenarios. Mini case study. At least 4 prompts requiring deep thought.

[CONSTRUCTIVE] 6. COMMON MISCONCEPTIONS
At least 5 common errors. For each: why students make it, then the correct understanding.

[INTERACTIVE] 7. COLLABORATIVE ACTIVITIES
Pair/group activities with specific prompts. Think-pair-share. Peer teaching exercise.

[INTERACTIVE] 8. REAL-WORLD APPLICATION
A substantial real-world scenario to solve collaboratively. Include reflection questions.

[PASSIVE] 9. SUMMARY & KEY TAKEAWAYS
Bullet-point recap of every major concept. A cheat sheet of key formulas/rules. Connection to next topic.

[CONSTRUCTIVE] 10. SELF-ASSESSMENT
5 self-check questions (recall, application, analysis). A "one-minute paper" prompt. Suggested further reading.

Formatting rules:
- Use **bold** for key terms
- Use bullet points with - for lists
- Start worked examples with "Example:"
- Write in clear, engaging academic language
- Aim for at least 2,500 words total
- Do NOT add any extra labels or tags beyond the [ICAP] ones shown above"""

        result = ask_groq(p, max_tokens=4000)

        code = d.get("classCode","").upper().strip()
        if code:
            init_db()
            conn = get_db()
            row = conn.execute("SELECT data FROM classes WHERE code=?", (code,)).fetchone()
            if row:
                cls = json.loads(row["data"])
                cls["notes"] = result
                conn.execute("UPDATE classes SET data=? WHERE code=?", (json.dumps(cls), code))
                conn.commit()
            conn.close()

        return jsonify({"success": True, "notes": result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_notes", methods=["POST"])
def get_notes():
    try:
        d = request.json
        code = d.get("classCode","").upper().strip()
        if not code:
            return jsonify({"success": False, "error": "No code"})
        init_db()
        conn = get_db()
        row = conn.execute("SELECT data FROM classes WHERE code=?", (code,)).fetchone()
        conn.close()
        if row:
            cls = json.loads(row["data"])
            return jsonify({"success": True, "notes": cls.get("notes","")})
        return jsonify({"success": False, "error": "Class not found"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  AI: SLIDESHOW DATA
# ══════════════════════════════════════════════════════════════
@app.route("/generate_slideshow_data", methods=["POST"])
def generate_slideshow_data():
    try:
        d = request.json
        topic = d.get("topic","Topic")
        level = d.get("level","Intermediate")
        duration = d.get("duration", 75)
        notes = d.get("notes","")
        language = d.get("language","English")

        context = f"Based on these detailed lecture notes:\n{notes[:3500]}" if notes and len(notes) > 100 else f"Topic: {topic}"

        p = f"""{context}

You are a university professor delivering a live {duration}-minute lecture on "{topic}" at {level} level in {language}.

Create 18 detailed lecture slides. Each slide should feel like you are ACTUALLY SPEAKING to students — the narration should be a full paragraph of spoken lecture content, not a summary. Think of it as a real transcript of what a brilliant professor would say.

RESPOND WITH ONLY A JSON ARRAY. No text before or after. No markdown fences.

Format per slide:
{{"title":"Slide title","bullets":["Point 1 — detailed sentence","Point 2 — detailed sentence","Point 3 — detailed sentence","Point 4 — detailed sentence"],"narration":"FULL spoken paragraph 80-150 words — conversational, engaging, like a real professor speaking. Include transitions, emphasis, examples, analogies. Say things like 'Now, what I want you to notice here is...', 'This is really important...', 'Think about it this way...', 'Let me give you a concrete example...'","icap":"passive","type":"content"}}

The 18 slides must cover:
1. Title slide — welcome and overview (type="title", icap="passive")
2. Why this matters — real-world hook, compelling opening story or statistic
3. Learning objectives — what students will be able to DO by the end
4. ICAP learning guide — briefly explain how today's session is structured
5. Core concept 1 — first major idea, deep explanation
6. Core concept 2 — second major idea
7. Core concept 3 — third major idea
8. Key definitions — precise terminology with examples
9. Worked example 1 — step by step (type="example", icap="active")
10. Worked example 2 — harder problem (type="example", icap="active")
11. Visual/Mental model — how to picture or remember the concepts
12. Common misconception 1 — what students get wrong and why (icap="constructive")
13. Common misconception 2 — another misconception
14. Critical thinking challenge — open question, push students to analyse (icap="constructive")
15. Peer discussion activity — pair/group exercise with prompts (type="activity", icap="interactive")
16. Real-world application — case study or industry example (icap="interactive")
17. Key takeaways — the 4-5 things students MUST remember (type="summary")
18. Closing — next steps, self-study advice, encouragement (type="title")

Rules:
- Narrations must be 80-150 words each — real spoken lecture voice
- Bullets must be complete informative sentences (not fragments)
- Include specific facts, numbers, examples relevant to {topic}
- Build logically — each slide should feel like it flows from the previous
- Level: {level} — calibrate depth accordingly
- Language: {language}

Return ONLY the JSON array starting with [ and ending with ]."""

        result = ask_groq(p, max_tokens=6000)
        slides = None

        for attempt in [
            lambda: json.loads(result),
            lambda: json.loads(re.search(r'\[[\s\S]*\]', result).group()),
            lambda: json.loads(re.sub(r'^```(?:json)?\s*|\s*```$', '', result.strip())),
        ]:
            try:
                slides = attempt()
                if slides and isinstance(slides, list) and len(slides) > 0:
                    break
            except:
                pass

        if slides:
            return jsonify({"success": True, "slides": slides})
        return jsonify({"success": False, "error": "AI did not return valid slides. Using local fallback."})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  AI: STUDY PLAN
# ══════════════════════════════════════════════════════════════
@app.route("/generate_study_plan", methods=["POST"])
def generate_study_plan():
    try:
        d = request.json
        topic = d.get("topic","")
        level = d.get("level","Intermediate")
        background = d.get("background","")
        language = d.get("language","English")

        p = f"""Create a personalised 7-day study plan for a student studying "{topic}" at {level} level.
Student background: {background if background else 'General learner'}.
Write in {language}.

For each day include: study focus (30-60 min), specific tasks, one self-check question, and a daily tip.
End with 3 recommended resources. Be specific and actionable."""

        result = ask_groq(p, max_tokens=1500)
        return jsonify({"success": True, "plan": result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  AI: LAYER 2 REAL-TIME TOOLS
# ══════════════════════════════════════════════════════════════
@app.route("/layer2/question", methods=["POST"])
def layer2_question():
    try:
        d = request.json
        p = f'Instructor. Topic:"{d.get("topic")}" Level:{d.get("level")}. Question:"{d.get("question")}"\nExplain clearly in 3 sentences. Note one common misconception. Give one follow-up question.'
        return jsonify({"result": ask_groq(p)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/layer2/confusion", methods=["POST"])
def layer2_confusion():
    try:
        d = request.json
        p = f'Instructor. Topic:"{d.get("topic")}" Level:{d.get("level")}. Confusion:"{d.get("confusion")}"\nGive a new analogy. Suggest a 3-minute rescue activity. End with one re-engage sentence.'
        return jsonify({"result": ask_groq(p)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/layer2/conceptcheck", methods=["POST"])
def layer2_conceptcheck():
    try:
        d = request.json
        p = f'Instructor. Topic:"{d.get("topic")}" Level:{d.get("level")}. Asked:"{d.get("question")}". {d.get("correct_pct")}% correct.\nInterpret this. What should the instructor do in the next 5 minutes?'
        return jsonify({"result": ask_groq(p)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/layer2/pacing", methods=["POST"])
def layer2_pacing():
    try:
        d = request.json
        total = int(d.get("total_duration", 75))
        elapsed = int(d.get("mins_elapsed", 0))
        remaining = total - elapsed
        p = f'Instructor. Topic:"{d.get("topic")}". {total}min total. {elapsed}min elapsed. {remaining}min remaining. On:"{d.get("current_segment")}".\nAre they on track? What to do now? What can be cut if needed?'
        return jsonify({"result": ask_groq(p)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/layer2/student_question", methods=["POST"])
def layer2_student_question():
    try:
        d = request.json
        p = f'Friendly tutor. Topic:"{d.get("topic")}". Student:{d.get("name")}, {d.get("year")}, background:{d.get("background")}, level:{d.get("level")}.\nQuestion:"{d.get("question")}"\nAnswer in under 150 words. Use plain language. End with encouragement using their name.'
        return jsonify({"result": ask_groq(p)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/generate_video_script", methods=["POST"])
def generate_video_script():
    try:
        d = request.json
        p = f'Write a 5-minute video lecture script for "{d.get("topic")}" at {d.get("level")} level. Include: hook intro, what it is, how it works, real example, common mistakes, summary outro. Use natural spoken language with [PAUSE] markers.'
        return jsonify({"result": ask_groq(p, max_tokens=1200)})
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/ai_feedback", methods=["POST"])
def ai_feedback():
    try:
        d = request.json
        p = f"""You are a helpful academic assistant. A student submitted an assignment.

Assignment: {d.get("title","")}
Instructions: {d.get("description","")}
Max score: {d.get("maxScore",100)}
Student submission: {d.get("content","")}

Give concise, encouraging feedback in 3-4 sentences. Note one strength and one area to improve. Do not assign a score."""
        result = ask_groq(p, max_tokens=300)
        return jsonify({"success": True, "feedback": result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  AI: QUIZ GENERATOR
# ══════════════════════════════════════════════════════════════
@app.route("/generate_quiz", methods=["POST"])
def generate_quiz():
    try:
        d = request.json
        topic = d.get("topic","")
        level = d.get("level","Intermediate")
        notes = d.get("notes","")
        language = d.get("language","English")

        context = f"Based on these lecture notes:\n{notes[:1200]}" if notes else f"Based on the topic: {topic}"
        p = f"""{context}

Create exactly 5 multiple choice questions for {level}-level students on "{topic}". Write in {language}.

Return ONLY valid JSON array, no preamble or markdown:
[
  {{"q":"Question?","options":["A","B","C","D"],"ans":0,"exp":"Why A is correct"}},
  ...
]

The "ans" field is the 0-indexed position of the correct answer."""

        result = ask_groq(p, max_tokens=900)
        match = re.search(r'\[.*\]', result, re.DOTALL)
        questions = json.loads(match.group() if match else result)
        return jsonify({"success": True, "questions": questions})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  REACTIONS
# ══════════════════════════════════════════════════════════════
@app.route("/save_reaction", methods=["POST"])
def save_reaction():
    try:
        d = request.json
        init_db()
        conn = get_db()
        conn.execute(
            """DELETE FROM reactions WHERE class_code=? AND student_email=?
            AND datetime(created_at) > datetime('now', '-30 seconds')""",
            (d.get("classCode",""), d.get("studentEmail",""))
        )
        conn.execute(
            """INSERT INTO reactions (id,class_code,student_email,student_name,reaction,created_at)
            VALUES(?,?,?,?,?,datetime('now'))""",
            (str(uuid.uuid4()), d.get("classCode",""), d.get("studentEmail",""),
             d.get("studentName",""), d.get("reaction",""))
        )
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_reactions", methods=["POST"])
def get_reactions():
    try:
        d = request.json
        code = d.get("classCode","")
        init_db()
        conn = get_db()
        rows = conn.execute(
            """SELECT reaction, COUNT(*) as cnt FROM reactions
            WHERE class_code=? AND datetime(created_at) > datetime('now', '-5 minutes')
            GROUP BY reaction""", (code,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "reactions": {r["reaction"]: r["cnt"] for r in rows}})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  ATTENDANCE
# ══════════════════════════════════════════════════════════════
@app.route("/save_attendance", methods=["POST"])
def save_attendance():
    try:
        d = request.json
        init_db()
        conn = get_db()
        existing = conn.execute(
            "SELECT id FROM attendance WHERE class_code=? AND student_name=? AND session_date=?",
            (d.get("classCode",""), d.get("studentName",""), d.get("sessionDate",""))
        ).fetchone()
        if existing:
            conn.execute(
                "UPDATE attendance SET present=? WHERE id=?",
                (1 if d.get("present", True) else 0, existing["id"])
            )
        else:
            conn.execute(
                """INSERT INTO attendance (id,class_code,teacher_email,student_name,session_date,present,created_at)
                VALUES(?,?,?,?,?,?,datetime('now'))""",
                (str(uuid.uuid4()), d.get("classCode",""), d.get("teacherEmail",""),
                 d.get("studentName",""), d.get("sessionDate",""),
                 1 if d.get("present", True) else 0)
            )
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/get_attendance", methods=["POST"])
def get_attendance():
    try:
        d = request.json
        code = d.get("classCode","")
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM attendance WHERE class_code=? ORDER BY session_date DESC, student_name",
            (code,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "attendance": [dict(r) for r in rows]})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  DISCUSSION BOARD
# ══════════════════════════════════════════════════════════════
@app.route("/discussion/post", methods=["POST"])
def discussion_post():
    try:
        d = request.json
        init_db()
        conn = get_db()
        did = str(uuid.uuid4())
        conn.execute(
            """INSERT INTO discussions (id,class_code,student_name,student_email,question,created_at)
            VALUES(?,?,?,?,?,datetime('now'))""",
            (did, d.get("classCode",""), d.get("studentName",""),
             d.get("studentEmail",""), d.get("question",""))
        )
        conn.commit()
        conn.close()
        return jsonify({"success": True, "id": did})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/discussion/get", methods=["POST"])
def discussion_get():
    try:
        d = request.json
        code = d.get("classCode","")
        init_db()
        conn = get_db()
        rows = conn.execute(
            "SELECT * FROM discussions WHERE class_code=? ORDER BY created_at DESC", (code,)
        ).fetchall()
        conn.close()
        return jsonify({"success": True, "posts": [dict(r) for r in rows]})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/discussion/reply", methods=["POST"])
def discussion_reply():
    try:
        d = request.json
        init_db()
        conn = get_db()
        conn.execute(
            "UPDATE discussions SET reply=?, replied_by=? WHERE id=?",
            (d.get("reply",""), d.get("repliedBy",""), d.get("id",""))
        )
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  POWERPOINT EXPORT
# ══════════════════════════════════════════════════════════════
@app.route("/generate_slides", methods=["POST"])
def generate_slides():
    try:
        d = request.json
        topic = d.get("topic","Topic")
        level = d.get("level","Intermediate")
        duration = d.get("duration", 75)
        objectives = d.get("objectives","")
        style = d.get("style","Lecture-based")
        notes = d.get("notes","")

        GREEN   = RGBColor(0x2d,0x6a,0x4f)
        WHITE   = RGBColor(0xFF,0xFF,0xFF)
        DARK    = RGBColor(0x1a,0x1a,0x1a)
        LGRAY   = RGBColor(0xF7,0xF5,0xF2)
        ACCENT  = RGBColor(0x74,0xC6,0x9D)
        DKGREEN = RGBColor(0x1B,0x43,0x32)
        SOFTW   = RGBColor(0xEC,0xEC,0xEC)

        ICAP_COLORS = {
            "PASSIVE":      RGBColor(0x6C,0x75,0x7D),
            "ACTIVE":       RGBColor(0x0D,0x6E,0xFD),
            "CONSTRUCTIVE": RGBColor(0xF4,0xA2,0x61),
            "INTERACTIVE":  RGBColor(0xE6,0x39,0x46),
        }
        ICAP_LABELS = {
            "PASSIVE":      "PASSIVE — Receiving",
            "ACTIVE":       "ACTIVE — Manipulating",
            "CONSTRUCTIVE": "CONSTRUCTIVE — Generating",
            "INTERACTIVE":  "INTERACTIVE — Dialoguing",
        }

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def rect(slide, l, t, w, h, c):
            s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
            return s

        def txt(slide, text, l, t, w, h, sz=18, bold=False, color=None, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size = Pt(sz); run.font.bold = bold
            run.font.color.rgb = color if color else DARK

        slide_sections = []
        if notes and len(notes) > 200:
            icap_tags    = re.findall(r'\[(PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]', notes)
            headers      = re.findall(r'\[(?:PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]\s*\d+\.\s*([^\n]+)', notes)
            sections     = re.split(r'\[(?:PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]\s*\d+\.', notes)
            for i, header in enumerate(headers[:12]):
                content = sections[i+1] if i+1 < len(sections) else ""
                lines = [re.sub(r'^[\s•\-\*\d\.]+','',l).strip() for l in content.split('\n')]
                lines = [l for l in lines if len(l) > 10][:5]
                tag   = icap_tags[i] if i < len(icap_tags) else "PASSIVE"
                if lines:
                    slide_sections.append({"title":header.strip(),"bullets":lines,"icap":tag})

        if len(slide_sections) < 5:
            try:
                gen_p = f"""For a {level}-level lecture on "{topic}", provide slide content.
Return ONLY valid JSON array:
[
  {{"title":"Learning Objectives","bullets":["Obj 1","Obj 2","Obj 3"],"icap":"PASSIVE"}},
  {{"title":"Why This Matters","bullets":["Point 1","Point 2","Point 3"],"icap":"PASSIVE"}},
  {{"title":"Core Concept","bullets":["Idea 1","Idea 2","Idea 3"],"icap":"ACTIVE"}},
  {{"title":"How It Works","bullets":["Step 1","Step 2","Step 3"],"icap":"ACTIVE"}},
  {{"title":"Worked Example","bullets":["Problem","Step 1","Result"],"icap":"ACTIVE"}},
  {{"title":"Common Misconceptions","bullets":["Error 1","Error 2","Error 3"],"icap":"CONSTRUCTIVE"}},
  {{"title":"Discussion Activity","bullets":["Prompt 1","Prompt 2","Prompt 3"],"icap":"INTERACTIVE"}},
  {{"title":"Real-World Applications","bullets":["App 1","App 2","App 3"],"icap":"CONSTRUCTIVE"}},
  {{"title":"Key Takeaways","bullets":["Key 1","Key 2","Key 3","Key 4"],"icap":"PASSIVE"}}
]
Make every bullet a complete, informative sentence about {topic}."""
                ai_result = ask_groq(gen_p, max_tokens=2000)
                match = re.search(r'\[.*\]', ai_result, re.DOTALL)
                if match:
                    slide_sections = json.loads(match.group())
            except:
                pass

        if len(slide_sections) < 3:
            objs = [o.strip() for o in objectives.split('\n') if o.strip()][:4]
            slide_sections = [
                {"title":"Learning Objectives",   "bullets":objs or [f"Understand {topic}","Apply key concepts","Analyse and evaluate","Connect to practice"],"icap":"PASSIVE"},
                {"title":"Why This Matters",       "bullets":[f"Real relevance of {topic}","Industry applications","What problem it solves","Why professionals need this"],"icap":"PASSIVE"},
                {"title":"Core Concepts",          "bullets":["Fundamental definitions","Key properties","How components relate","Underlying logic"],"icap":"ACTIVE"},
                {"title":"Worked Example",         "bullets":["Define the problem","Choose the approach","Apply step by step","Interpret the result"],"icap":"ACTIVE"},
                {"title":"Common Misconceptions",  "bullets":["Confusing similar concepts","Skipping assumptions","Over-generalising","Ignoring edge cases"],"icap":"CONSTRUCTIVE"},
                {"title":"Discussion Activity",    "bullets":["Discuss: What is the key insight?","Create your own example","Identify one point of confusion","Prepare to share findings"],"icap":"INTERACTIVE"},
                {"title":"Real-World Application", "bullets":[f"{topic} in industry","A recent case study","Theory meets professional practice","What experts wish they learned earlier"],"icap":"CONSTRUCTIVE"},
                {"title":"Key Takeaways",          "bullets":["Core definition and importance","Start with intuition, then formalise","Practice with varied examples","Connect to the bigger picture"],"icap":"PASSIVE"},
            ]

        s1 = prs.slides.add_slide(blank)
        rect(s1, 0, 0, 13.33, 7.5, GREEN)
        rect(s1, 0, 5.5, 13.33, 2.0, DKGREEN)
        rect(s1, 0.5, 4.7, 12.33, 0.05, ACCENT)
        txt(s1, "LectureAI", 0.5, 0.4, 12, 0.5, sz=13, color=ACCENT, align=PP_ALIGN.CENTER)
        txt(s1, topic, 0.5, 1.1, 12, 2.8, sz=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s1, f"{level}  ·  {duration} min  ·  {style}", 0.5, 3.5, 12, 0.7, sz=18, color=ACCENT, align=PP_ALIGN.CENTER)
        txt(s1, "Human-AI Co-Orchestration in Education", 0.5, 5.9, 12, 0.5, sz=13, color=SOFTW, align=PP_ALIGN.CENTER)
        txt(s1, "Powered by ICAP Framework  ·  Chi & Wylie (2014)", 0.5, 6.5, 12, 0.6, sz=11, color=ACCENT, align=PP_ALIGN.CENTER)

        s2 = prs.slides.add_slide(blank)
        rect(s2, 0, 0, 13.33, 7.5, LGRAY)
        rect(s2, 0, 0, 13.33, 1.4, GREEN)
        txt(s2, "ICAP Framework: How This Lecture Is Designed", 0.4, 0.2, 12.5, 1.0, sz=26, bold=True, color=WHITE)
        txt(s2, "LectureAI", 11.5, 0.22, 1.5, 0.5, sz=10, color=ACCENT)
        for idx, (label, desc, ic) in enumerate([
            ("PASSIVE",      "Receiving\nListening, reading", ICAP_COLORS["PASSIVE"]),
            ("ACTIVE",       "Manipulating\nHighlighting, repeating", ICAP_COLORS["ACTIVE"]),
            ("CONSTRUCTIVE", "Generating\nExplaining, creating", ICAP_COLORS["CONSTRUCTIVE"]),
            ("INTERACTIVE",  "Dialoguing\nDebating, co-creating", ICAP_COLORS["INTERACTIVE"]),
        ]):
            bx = 0.5 + idx * 3.1
            rect(s2, bx, 1.8, 2.9, 0.12, ic)
            txt(s2, label, bx+0.1, 2.1, 2.7, 0.6, sz=20, bold=True, color=DARK)
            for li, line in enumerate(desc.split('\n')):
                txt(s2, line, bx+0.1, 2.8+li*0.4, 2.7, 0.4, sz=13, color=DARK)
        txt(s2, "Higher engagement levels → deeper learning outcomes (Chi & Wylie, 2014)",
            0.5, 5.2, 12.33, 0.5, sz=14, bold=True, color=GREEN)

        for idx, section in enumerate(slide_sections):
            s = prs.slides.add_slide(blank)
            is_dark = idx % 2 == 1
            bg = DKGREEN if is_dark else LGRAY
            text_c = WHITE if is_dark else DARK
            rect(s, 0, 0, 13.33, 7.5, bg)
            rect(s, 0, 0, 13.33, 1.4, GREEN)
            title_text = section.get("title","")
            icap_tag   = str(section.get("icap","PASSIVE")).upper()
            icap_color = ICAP_COLORS.get(icap_tag, ACCENT)
            icap_label = ICAP_LABELS.get(icap_tag, icap_tag)
            txt(s, title_text, 0.4, 0.2, 10.5, 1.0, sz=26, bold=True, color=WHITE)
            txt(s, "LectureAI", 11.5, 0.22, 1.5, 0.5, sz=10, color=ACCENT)
            rect(s, 0.4, 1.5, 3.2, 0.35, icap_color)
            txt(s, icap_label, 0.5, 1.52, 3.0, 0.3, sz=11, bold=True, color=WHITE)
            txt(s, f"Slide {idx+3}", 12.0, 1.52, 1.0, 0.3, sz=10, color=ACCENT if is_dark else GREEN)
            bullets = section.get("bullets", [])[:6]
            y = 2.15
            for b in bullets:
                rect(s, 0.5, y+0.07, 0.06, 0.32, ACCENT)
                txt(s, str(b), 0.75, y, 12.0, 0.5, sz=15, color=text_c)
                y += 0.62

        sc = prs.slides.add_slide(blank)
        rect(sc, 0, 0, 13.33, 7.5, DKGREEN)
        rect(sc, 0.5, 3.0, 12.33, 0.05, ACCENT)
        txt(sc, "Thank You", 0.5, 1.2, 12, 1.5, sz=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sc, f"Questions about {topic}?", 0.5, 3.2, 12, 0.8, sz=22, color=ACCENT, align=PP_ALIGN.CENTER)
        txt(sc, "The best way to learn is to explain it to someone else.", 0.5, 4.4, 12, 0.6, sz=14, color=SOFTW, align=PP_ALIGN.CENTER)
        txt(sc, "Built with LectureAI  ·  Human-AI Co-Orchestration  ·  ICAP Framework", 0.5, 6.5, 12, 0.6, sz=11, color=ACCENT, align=PP_ALIGN.CENTER)

        buf = io.BytesIO()
        prs.save(buf); buf.seek(0)
        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"LectureAI_{topic[:30].replace(' ','_')}.pptx"
        )
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"success": False, "error": str(e)})

# ══════════════════════════════════════════════════════════════
#  LECTURE LIBRARY
# ══════════════════════════════════════════════════════════════
@app.route("/library/save", methods=["POST"])
def library_save():
    try:
        d = request.json
        if not d.get("notes") or not d.get("topic"):
            return jsonify({"success": False, "error": "Topic and notes are required"})
        init_db()
        conn = get_db()
        existing = conn.execute(
            "SELECT id FROM lecture_library WHERE teacher_email=? AND topic=?",
            (d.get("teacherEmail",""), d.get("topic",""))
        ).fetchone()
        lid = existing["id"] if existing else str(uuid.uuid4())
        year = str(__import__('datetime').datetime.now().year)
        if existing:
            conn.execute("""UPDATE lecture_library SET
                title=?, notes=?, subject=?, level=?, institution=?,
                class_code=?, is_public=?, saved_at=datetime('now'), year=?, teacher_name=?
                WHERE id=?""",
                (d.get("title", d.get("topic","")), d.get("notes",""), d.get("subject",""),
                 d.get("level","Intermediate"), d.get("institution",""), d.get("classCode",""),
                 1 if d.get("isPublic", True) else 0, year, d.get("teacherName",""), lid))
        else:
            conn.execute("""INSERT INTO lecture_library
                (id, teacher_email, teacher_name, title, topic, subject, level,
                 institution, notes, class_code, is_public, view_count, saved_at, year)
                VALUES(?,?,?,?,?,?,?,?,?,?,?,0,datetime('now'),?)""",
                (lid, d.get("teacherEmail",""), d.get("teacherName",""),
                 d.get("title", d.get("topic","")), d.get("topic",""), d.get("subject",""),
                 d.get("level","Intermediate"), d.get("institution",""), d.get("notes",""),
                 d.get("classCode",""), 1 if d.get("isPublic", True) else 0, year))
        conn.commit()
        conn.close()
        return jsonify({"success": True, "id": lid, "updated": bool(existing)})
    except Exception as e:
        print("library_save error:", e)
        return jsonify({"success": False, "error": str(e)})

@app.route("/library/list", methods=["POST"])
def library_list():
    try:
        d = request.json or {}
        search = d.get("search","").strip()
        subject = d.get("subject","").strip()
        level = d.get("level","").strip()
        year = d.get("year","").strip()
        teacher_email = d.get("teacherEmail","").strip()
        page = int(d.get("page", 1))
        per_page = 20
        init_db()
        conn = get_db()
        query = "SELECT id, teacher_name, teacher_email, title, topic, subject, level, institution, year, view_count, saved_at FROM lecture_library WHERE is_public=1"
        params = []
        if search:
            query += " AND (topic LIKE ? OR title LIKE ? OR subject LIKE ? OR teacher_name LIKE ?)"
            s = f"%{search}%"
            params += [s, s, s, s]
        if subject:
            query += " AND subject LIKE ?"
            params.append(f"%{subject}%")
        if level:
            query += " AND level=?"
            params.append(level)
        if year:
            query += " AND year=?"
            params.append(year)
        if teacher_email:
            query += " AND teacher_email=?"
            params.append(teacher_email)
        count_row = conn.execute(f"SELECT COUNT(*) as cnt FROM ({query})", params).fetchone()
        total = count_row["cnt"] if count_row else 0
        query += " ORDER BY saved_at DESC LIMIT ? OFFSET ?"
        params += [per_page, (page-1)*per_page]
        rows = conn.execute(query, params).fetchall()
        conn.close()
        return jsonify({
            "success": True,
            "lectures": [dict(r) for r in rows],
            "total": total,
            "page": page,
            "pages": max(1, -(-total // per_page))
        })
    except Exception as e:
        print("library_list error:", e)
        return jsonify({"success": False, "error": str(e)})

@app.route("/library/get", methods=["POST"])
def library_get():
    try:
        d = request.json
        lid = d.get("id","")
        if not lid:
            return jsonify({"success": False, "error": "No ID provided"})
        init_db()
        conn = get_db()
        row = conn.execute("SELECT * FROM lecture_library WHERE id=?", (lid,)).fetchone()
        if row:
            conn.execute("UPDATE lecture_library SET view_count=view_count+1 WHERE id=?", (lid,))
            conn.commit()
        conn.close()
        if row:
            return jsonify({"success": True, "lecture": dict(row)})
        return jsonify({"success": False, "error": "Lecture not found"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/library/delete", methods=["POST"])
def library_delete():
    try:
        d = request.json
        lid = d.get("id","")
        email = d.get("teacherEmail","")
        init_db()
        conn = get_db()
        conn.execute("DELETE FROM lecture_library WHERE id=? AND teacher_email=?", (lid, email))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/library/subjects", methods=["POST"])
def library_subjects():
    try:
        init_db()
        conn = get_db()
        subjects = conn.execute(
            "SELECT DISTINCT subject FROM lecture_library WHERE is_public=1 AND subject!='' ORDER BY subject"
        ).fetchall()
        years = conn.execute(
            "SELECT DISTINCT year FROM lecture_library WHERE is_public=1 AND year!='' ORDER BY year DESC"
        ).fetchall()
        levels = conn.execute(
            "SELECT DISTINCT level FROM lecture_library WHERE is_public=1 AND level!='' ORDER BY level"
        ).fetchall()
        conn.close()
        return jsonify({
            "success": True,
            "subjects": [r["subject"] for r in subjects],
            "years": [r["year"] for r in years],
            "levels": [r["level"] for r in levels],
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
