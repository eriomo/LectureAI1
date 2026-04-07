import io
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import services.ai_service as ai


GREEN   = RGBColor(0x2d, 0x6a, 0x4f)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1a, 0x1a, 0x1a)
LGRAY   = RGBColor(0xF7, 0xF5, 0xF2)
ACCENT  = RGBColor(0x74, 0xC6, 0x9D)
DKGREEN = RGBColor(0x1B, 0x43, 0x32)
SOFTW   = RGBColor(0xEC, 0xEC, 0xEC)

ICAP_COLORS = {
    "PASSIVE":      RGBColor(0x6C, 0x75, 0x7D),
    "ACTIVE":       RGBColor(0x0D, 0x6E, 0xFD),
    "CONSTRUCTIVE": RGBColor(0xF4, 0xA2, 0x61),
    "INTERACTIVE":  RGBColor(0xE6, 0x39, 0x46),
}
ICAP_LABELS = {
    "PASSIVE":      "PASSIVE — Receiving",
    "ACTIVE":       "ACTIVE — Manipulating",
    "CONSTRUCTIVE": "CONSTRUCTIVE — Generating",
    "INTERACTIVE":  "INTERACTIVE — Dialoguing",
}


def _rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s


def _txt(slide, text, l, t, w, h, sz=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = color if color else DARK


def _parse_sections(notes: str, topic: str, level: str) -> list:
    """Extract slide sections from notes or generate via AI."""
    slide_sections = []
    if notes and len(notes) > 200:
        icap_tags  = re.findall(r'\[(PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]', notes)
        headers    = re.findall(r'\[(?:PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]\s*\d+\.\s*([^\n]+)', notes)
        sections   = re.split(r'\[(?:PASSIVE|ACTIVE|CONSTRUCTIVE|INTERACTIVE)\]\s*\d+\.', notes)
        for i, header in enumerate(headers[:12]):
            content = sections[i + 1] if i + 1 < len(sections) else ""
            lines = [re.sub(r'^[\s•\-\*\d\.]+', '', l).strip() for l in content.split('\n')]
            lines = [l for l in lines if len(l) > 10][:5]
            tag   = icap_tags[i] if i < len(icap_tags) else "PASSIVE"
            if lines:
                slide_sections.append({"title": header.strip(), "bullets": lines, "icap": tag})

    if len(slide_sections) < 5:
        try:
            gen_p = f"""For a {level}-level lecture on "{topic}", generate slide content.
Return ONLY valid JSON array (no markdown):
[
  {{"title":"Learning Objectives","bullets":["Obj 1","Obj 2","Obj 3"],"icap":"PASSIVE"}},
  {{"title":"Why This Matters","bullets":["Point 1","Point 2","Point 3"],"icap":"PASSIVE"}},
  {{"title":"Core Concept","bullets":["Idea 1","Idea 2","Idea 3"],"icap":"ACTIVE"}},
  {{"title":"How It Works","bullets":["Step 1","Step 2","Step 3"],"icap":"ACTIVE"}},
  {{"title":"Worked Example","bullets":["Problem","Step 1","Result"],"icap":"ACTIVE"}},
  {{"title":"Common Misconceptions","bullets":["Error 1","Error 2","Error 3"],"icap":"CONSTRUCTIVE"}},
  {{"title":"Discussion Activity","bullets":["Prompt 1","Prompt 2","Prompt 3"],"icap":"INTERACTIVE"}},
  {{"title":"Real-World Applications","bullets":["App 1","App 2","App 3"],"icap":"CONSTRUCTIVE"}},
  {{"title":"Key Takeaways","bullets":["Key 1","Key 2","Key 3","Key 4"],"icap":"PASSIVE"}}
]
Make every bullet a complete, informative sentence about {topic}."""
            raw = ai.ask(gen_p, max_tokens=2000)
            match = re.search(r'\[.*\]', raw, re.DOTALL)
            if match:
                slide_sections = json.loads(match.group())
        except Exception:
            pass

    if len(slide_sections) < 3:
        slide_sections = [
            {"title": "Learning Objectives",    "bullets": [f"Understand {topic}", "Apply key concepts", "Analyse and evaluate", "Connect to practice"], "icap": "PASSIVE"},
            {"title": "Why This Matters",        "bullets": [f"Real relevance of {topic}", "Industry applications", "What problem it solves", "Why professionals need this"], "icap": "PASSIVE"},
            {"title": "Core Concepts",           "bullets": ["Fundamental definitions", "Key properties", "How components relate", "Underlying logic"], "icap": "ACTIVE"},
            {"title": "Worked Example",          "bullets": ["Define the problem", "Choose the approach", "Apply step by step", "Interpret the result"], "icap": "ACTIVE"},
            {"title": "Common Misconceptions",   "bullets": ["Confusing similar concepts", "Skipping assumptions", "Over-generalising", "Ignoring edge cases"], "icap": "CONSTRUCTIVE"},
            {"title": "Discussion Activity",     "bullets": ["Discuss the key insight", "Create your own example", "Identify one confusion", "Prepare to share"], "icap": "INTERACTIVE"},
            {"title": "Key Takeaways",           "bullets": ["Core definition", "Start with intuition", "Practice with varied examples", "Connect to bigger picture"], "icap": "PASSIVE"},
        ]
    return slide_sections


def build_pptx(d: dict) -> io.BytesIO:
    topic      = d.get("topic", "Topic")
    level      = d.get("level", "Intermediate")
    duration   = d.get("duration", 75)
    objectives = d.get("objectives", "")
    style      = d.get("style", "Lecture-based")
    notes      = d.get("notes", "")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Title slide ──────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _rect(s1, 0, 0, 13.33, 7.5, GREEN)
    _rect(s1, 0, 5.5, 13.33, 2.0, DKGREEN)
    _rect(s1, 0.5, 4.7, 12.33, 0.05, ACCENT)
    _txt(s1, "LectureAI", 0.5, 0.4, 12, 0.5, sz=13, color=ACCENT, align=PP_ALIGN.CENTER)
    _txt(s1, topic, 0.5, 1.1, 12, 2.8, sz=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s1, f"{level}  ·  {duration} min  ·  {style}", 0.5, 3.5, 12, 0.7, sz=18, color=ACCENT, align=PP_ALIGN.CENTER)
    _txt(s1, "Human-AI Co-Orchestration in Education", 0.5, 5.9, 12, 0.5, sz=13, color=SOFTW, align=PP_ALIGN.CENTER)
    _txt(s1, "Powered by ICAP Framework  ·  Chi & Wylie (2014)", 0.5, 6.5, 12, 0.5, sz=11, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── ICAP overview slide ──────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _rect(s2, 0, 0, 13.33, 7.5, LGRAY)
    _rect(s2, 0, 0, 13.33, 1.4, GREEN)
    _txt(s2, "ICAP Framework: How This Lecture Is Designed", 0.4, 0.2, 12.5, 1.0, sz=26, bold=True, color=WHITE)
    _txt(s2, "LectureAI", 11.5, 0.22, 1.5, 0.5, sz=10, color=ACCENT)
    for idx, (label, desc, ic) in enumerate([
        ("PASSIVE",      "Receiving\nListening, reading",    ICAP_COLORS["PASSIVE"]),
        ("ACTIVE",       "Manipulating\nHighlighting, doing", ICAP_COLORS["ACTIVE"]),
        ("CONSTRUCTIVE", "Generating\nExplaining, creating",  ICAP_COLORS["CONSTRUCTIVE"]),
        ("INTERACTIVE",  "Dialoguing\nDebating, co-creating", ICAP_COLORS["INTERACTIVE"]),
    ]):
        bx = 0.5 + idx * 3.1
        _rect(s2, bx, 1.8, 2.9, 0.12, ic)
        _txt(s2, label, bx + 0.1, 2.1, 2.7, 0.6, sz=20, bold=True, color=DARK)
        for li, line in enumerate(desc.split('\n')):
            _txt(s2, line, bx + 0.1, 2.8 + li * 0.4, 2.7, 0.4, sz=13, color=DARK)
    _txt(s2, "Higher engagement → deeper learning outcomes (Chi & Wylie, 2014)", 0.5, 5.2, 12.33, 0.5, sz=14, bold=True, color=GREEN)

    # ── Content slides ───────────────────────────────────────────
    slide_sections = _parse_sections(notes, topic, level)
    for idx, section in enumerate(slide_sections):
        s = prs.slides.add_slide(blank)
        is_dark  = idx % 2 == 1
        bg       = DKGREEN if is_dark else LGRAY
        text_c   = WHITE if is_dark else DARK
        _rect(s, 0, 0, 13.33, 7.5, bg)
        _rect(s, 0, 0, 13.33, 1.4, GREEN)
        title_text = section.get("title", "")
        icap_tag   = str(section.get("icap", "PASSIVE")).upper()
        icap_color = ICAP_COLORS.get(icap_tag, ACCENT)
        icap_label = ICAP_LABELS.get(icap_tag, icap_tag)
        _txt(s, title_text, 0.4, 0.2, 10.5, 1.0, sz=26, bold=True, color=WHITE)
        _txt(s, "LectureAI", 11.5, 0.22, 1.5, 0.5, sz=10, color=ACCENT)
        _rect(s, 0.4, 1.5, 3.2, 0.35, icap_color)
        _txt(s, icap_label, 0.5, 1.52, 3.0, 0.3, sz=11, bold=True, color=WHITE)
        _txt(s, f"Slide {idx + 3}", 12.0, 1.52, 1.0, 0.3, sz=10, color=ACCENT if is_dark else GREEN)
        bullets = section.get("bullets", [])[:6]
        y = 2.15
        for b in bullets:
            _rect(s, 0.5, y + 0.07, 0.06, 0.32, ACCENT)
            _txt(s, str(b), 0.75, y, 12.0, 0.5, sz=15, color=text_c)
            y += 0.62

    # ── Thank you slide ──────────────────────────────────────────
    sc = prs.slides.add_slide(blank)
    _rect(sc, 0, 0, 13.33, 7.5, DKGREEN)
    _rect(sc, 0.5, 3.0, 12.33, 0.05, ACCENT)
    _txt(sc, "Thank You", 0.5, 1.2, 12, 1.5, sz=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sc, f"Questions about {topic}?", 0.5, 3.2, 12, 0.8, sz=22, color=ACCENT, align=PP_ALIGN.CENTER)
    _txt(sc, "The best way to learn is to explain it to someone else.", 0.5, 4.4, 12, 0.6, sz=14, color=SOFTW, align=PP_ALIGN.CENTER)
    _txt(sc, "Built with LectureAI  ·  Human-AI Co-Orchestration  ·  ICAP Framework", 0.5, 6.5, 12, 0.6, sz=11, color=ACCENT, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
